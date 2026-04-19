"""PowerPoint (.pptx) parser — converts slide decks to the Chapter IR."""

from __future__ import annotations

import tempfile
from pathlib import Path

from ..model import (
    Author,
    Chapter,
    Figure,
    InlineRun,
    ListBlock,
    ListItem,
    Paragraph,
    Section,
    Table,
    TableCell,
)


def parse_pptx(path: Path) -> Chapter:
    """Parse a .pptx file into a Chapter IR."""
    from pptx import Presentation

    prs = Presentation(str(path))
    chapter = Chapter()

    slides_data = [_parse_slide(slide, chapter, idx) for idx, slide in enumerate(prs.slides)]

    # First slide is often a title slide
    if slides_data and slides_data[0].get("is_title_slide"):
        td = slides_data[0]
        chapter.title = td.get("title", "")
        if td.get("authors"):
            chapter.authors = td["authors"]
        slides_data = slides_data[1:]
    elif slides_data:
        chapter.title = slides_data[0].get("title", "")

    # Convert remaining slides to sections
    chapter.sections = _slides_to_sections(slides_data)

    return chapter


# ---------------------------------------------------------------------------
# Slide parsing
# ---------------------------------------------------------------------------


def _parse_slide(slide, chapter: Chapter, idx: int) -> dict:
    """Parse a single slide into a dict of extracted content."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    data: dict = {
        "title": "",
        "content": [],      # list of Block objects
        "is_title_slide": False,
        "notes": "",
        "authors": [],
    }

    # Sort shapes by position: top then left
    shapes = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))

    # Find the title
    if slide.shapes.title and slide.shapes.title.text.strip():
        data["title"] = slide.shapes.title.text.strip()

    # Detect title slide: has a subtitle placeholder or is the first slide with
    # a layout name containing "title"
    layout_name = slide.slide_layout.name.lower() if slide.slide_layout else ""
    is_title_layout = "title slide" in layout_name or "title only" in layout_name

    subtitle_text = ""
    for shape in shapes:
        ph_idx = _placeholder_idx(shape)
        if ph_idx == 1 and "title slide" in layout_name:
            subtitle_text = shape.text_frame.text if shape.has_text_frame else ""

    if is_title_layout and idx == 0 and data["title"]:
        data["is_title_slide"] = True
        if subtitle_text:
            data["authors"] = _parse_subtitle_authors(subtitle_text)

    # Extract content from non-title shapes
    for shape in shapes:
        # Skip the title shape itself
        if shape == slide.shapes.title:
            continue

        # Skip subtitle on title slides (already processed)
        if data["is_title_slide"] and _placeholder_idx(shape) == 1:
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = _parse_table(shape.table)
            if table:
                data["content"].append(table)
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            fig = _extract_image(shape, chapter)
            if fig:
                data["content"].append(fig)
            continue

        if shape.has_text_frame:
            # If title is empty, try to get it from the first text box
            if not data["title"] and _looks_like_title(shape):
                data["title"] = shape.text_frame.text.strip()
                continue

            blocks = _parse_text_frame(shape.text_frame)
            data["content"].extend(blocks)

    # Speaker notes
    try:
        notes = slide.notes_slide.notes_text_frame.text
        if notes.strip():
            data["notes"] = notes.strip()
    except Exception:
        pass

    return data


def _looks_like_title(shape) -> bool:
    """Heuristic: a text box near the top with short text is likely a title."""
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    if not text or len(text) > 100:
        return False
    # Near the top of the slide
    if shape.top is not None and shape.top < 500000:  # ~0.5 inches
        return True
    return False


# ---------------------------------------------------------------------------
# Text frame -> blocks
# ---------------------------------------------------------------------------


def _parse_text_frame(tf) -> list:
    """Convert a text frame's paragraphs into Block objects."""
    blocks: list = []
    pending_list: list[tuple[list[InlineRun], int]] = []  # (runs, level)

    for para in tf.paragraphs:
        runs = _parse_runs(para)
        text = "".join(r.text for r in runs).strip()
        if not text:
            if pending_list:
                blocks.append(_flush_bullet_list(pending_list))
                pending_list = []
            continue

        # Detect bullet lists: paragraphs with level >= 0 in content placeholders
        # In PPTX, most body text is at level 0 with bullets implied.
        # We treat multi-paragraph text frames as bullet lists.
        if len(tf.paragraphs) > 1 and _is_bullet_context(tf):
            pending_list.append((runs, para.level))
        else:
            if pending_list:
                blocks.append(_flush_bullet_list(pending_list))
                pending_list = []
            blocks.append(Paragraph(runs=runs))

    if pending_list:
        blocks.append(_flush_bullet_list(pending_list))

    return blocks


def _is_bullet_context(tf) -> bool:
    """Check if a text frame looks like a bullet list context."""
    # If more than one non-empty paragraph, treat as list
    non_empty = [p for p in tf.paragraphs if p.text.strip()]
    return len(non_empty) > 1


def _flush_bullet_list(pending: list[tuple[list[InlineRun], int]]) -> ListBlock:
    """Convert pending bullet items into a ListBlock."""
    items = [ListItem(runs=runs) for runs, _ in pending]
    pending.clear()
    return ListBlock(ordered=False, items=items)


def _parse_runs(para) -> list[InlineRun]:
    """Convert a PPTX paragraph's runs into InlineRun objects."""
    runs: list[InlineRun] = []
    for run in para.runs:
        text = run.text
        if not text:
            continue
        bold = bool(run.font.bold)
        italic = bool(run.font.italic)
        runs.append(InlineRun(text=text, bold=bold, italic=italic))
    return runs


# ---------------------------------------------------------------------------
# Table parsing
# ---------------------------------------------------------------------------


def _parse_table(tbl) -> Table | None:
    """Parse a PPTX table shape into a Table IR object."""
    rows_data: list[list[str]] = []
    for row in tbl.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows_data.append(cells)

    if not rows_data:
        return None

    headers = [TableCell(text=c) for c in rows_data[0]]
    body_rows = [[TableCell(text=c) for c in r] for r in rows_data[1:]]
    return Table(headers=headers, rows=body_rows)


# ---------------------------------------------------------------------------
# Image extraction
# ---------------------------------------------------------------------------


def _extract_image(shape, chapter: Chapter) -> Figure | None:
    """Extract an image from a picture shape."""
    try:
        img = shape.image
        blob = img.blob
        ext = img.ext or "png"
    except Exception:
        return None

    name = f"slide_image.{ext}"
    tmp_dir = Path(tempfile.mkdtemp(prefix="occonvert_"))
    tmp_path = tmp_dir / name
    tmp_path.write_bytes(blob)

    fig = Figure(
        source_path=str(tmp_path),
        output_filename="",
        caption="",
        alt_text=name,
    )
    chapter.images.append(fig)
    return fig


# ---------------------------------------------------------------------------
# Author parsing from subtitle
# ---------------------------------------------------------------------------


def _parse_subtitle_authors(text: str) -> list[Author]:
    """Parse subtitle text into Author objects.

    Common patterns:
    - "Jane Smith\nCarnegie Mellon University"
    - "Jane Smith, Carnegie Mellon University"
    """
    lines = [l.strip() for l in text.split("\n") if l.strip()]
    if not lines:
        return []

    # First line is name, subsequent lines are affiliation
    name = lines[0]
    institution = ", ".join(lines[1:]) if len(lines) > 1 else ""

    # If name contains comma, it might be "Name, University"
    if not institution and "," in name:
        parts = [p.strip() for p in name.split(",", 1)]
        name = parts[0]
        institution = parts[1]

    name_parts = name.rsplit(" ", 1)
    first = name_parts[0] if len(name_parts) > 1 else name
    last = name_parts[1] if len(name_parts) > 1 else ""

    return [Author(first=first, last=last, institution=institution)]


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------


def _placeholder_idx(shape) -> int | None:
    """Return the placeholder index of a shape, or None if not a placeholder."""
    try:
        return shape.placeholder_format.idx
    except (ValueError, AttributeError):
        return None


# ---------------------------------------------------------------------------
# Slides -> Sections
# ---------------------------------------------------------------------------


def _slides_to_sections(slides: list[dict]) -> list[Section]:
    """Convert slide data dicts into a flat list of sections."""
    sections: list[Section] = []

    for sd in slides:
        title = sd.get("title", "") or "Untitled Slide"
        content = list(sd.get("content", []))
        notes = sd.get("notes", "")

        # Add speaker notes as a comment paragraph
        if notes:
            comment_runs = [InlineRun(text=f"% NOTE: {notes}")]
            content.append(Paragraph(runs=comment_runs))

        # Add TODO comment for expanding slide content
        content.append(
            Paragraph(runs=[InlineRun(text="% TODO: Expand slide content into full prose")])
        )

        sections.append(Section(title=title, level=1, content=content))

    return sections
